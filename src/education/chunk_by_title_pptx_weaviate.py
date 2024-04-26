import glob
import os
from pptx import Presentation
import weaviate
from weaviate.config import AdditionalConfig
from weaviate.classes.config import Configure, DataType, Property
from unstructured.partition.pptx import partition_pptx
import concurrent.futures


def extract_text(file_name: str):
    prs = Presentation(file_name)
    elements = partition_pptx(
        filename=file_name,
        multipage_sections=True,
        infer_table_structure=True,
        include_page_breaks=True,
    )

    result_list = []
    for i, slide in enumerate(prs.slides):
        slide_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text += shape.text_frame.text + " "

        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text

        slide_elements = elements[i] if i < len(elements) else []
        element_text = []
        if isinstance(slide_elements, list):
            for element in slide_elements:
                if hasattr(element, "text"):
                    element_text.append(element.text)
                else:
                    element_text.append(str(element))
        else:
            element_text.append(str(slide_elements))

        combined_text = (
            f"{slide_text.strip()}\n\n{notes_text.strip()}\n\n{' '.join(element_text)}"
        )
        slide_title = slide_text.strip().split("\n")[0] if slide_text.strip() else ""
        base_name = os.path.splitext(os.path.basename(file_name))[0]

        result_list.append(
            {
                "title": slide_title,
                "content": combined_text,
                "source": f"{base_name} - Slide {i+1}",
            }
        )

    return result_list


def process_pptx(file_path):
    contents = extract_text(file_path)
    water_collection = w_client.collections.get(name="Water_pptx")
    for item in contents:
        water_collection.data.insert(item)


w_client = weaviate.connect_to_local(
    host="localhost", additional_config=AdditionalConfig(timeout=(600, 800))
)
try:
    collection = w_client.collections.create(
        name="Water_pptx",
        properties=[
            Property(name="title", data_type=DataType.TEXT),
            Property(name="content", data_type=DataType.TEXT),
            Property(name="source", data_type=DataType.TEXT),
        ],
        vectorizer_config=[
            Configure.NamedVectors.text2vec_transformers(
                name="title", source_properties=["title"]
            ),
            Configure.NamedVectors.text2vec_transformers(
                name="content", source_properties=["content"]
            ),
        ],
    )

    directory = "test"
    pptx_files = glob.glob(os.path.join(directory, "*.pptx"))
    with concurrent.futures.ProcessPoolExecutor(max_workers=6) as executor:
        executor.map(process_pptx, pptx_files)

    print("Data inserted successfully")

finally:
    w_client.close()
