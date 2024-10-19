import glob
import os
from datetime import datetime
import psycopg2
from psycopg2 import sql
from pptx import Presentation
from unstructured.partition.pptx import partition_pptx
import concurrent.futures
import pickle


def extract_text(file_name: str):
    prs = Presentation(file_name)
    elements = partition_pptx(
        filename=file_name,
        multipage_sections=True,
        infer_table_structure=True,
        include_page_breaks=True,
    )

    result_list = []
    for i, slide in enumerate(prs.slides):
        slide_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text += shape.text_frame.text + " "

        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text

        slide_elements = elements[i] if i < len(elements) else []
        element_text = []
        if isinstance(slide_elements, list):
            for element in slide_elements:
                if hasattr(element, "text"):
                    element_text.append(element.text)
                else:
                    element_text.append(str(element))
        else:
            element_text.append(str(slide_elements))

        combined_text = f"Page: {i+1}\n{slide_text.strip()}\n\n{notes_text.strip()}\n\n{' '.join(element_text)}"

        result_list.append(combined_text)

    return result_list


def process_pptx(file_path):
    record_id = os.path.splitext(os.path.basename(file_path))[0]

    text_list = extract_text(file_path)

    with open("processed_docs/ali_pickle/" + record_id + ".pptx" + ".pkl", "wb") as f:
        pickle.dump(text_list, f)

    text_str = "\n----------\n".join(map(str, text_list))

    with open("processed_docs/ali_txt/" + record_id + ".pptx" + ".txt", "w") as f:
        f.write(text_str)


conn = psycopg2.connect(
    database=os.getenv("POSTGRES_DB"),
    user=os.getenv("POSTGRES_USER"),
    password=os.getenv("POSTGRES_PASSWORD"),
    host=os.getenv("POSTGRES_HOST"),
    port=os.getenv("POSTGRES_PORT"),
)
cur = conn.cursor()
cur.execute("SELECT id FROM ali WHERE unstructure_time is NULL")
rows = cur.fetchall()

directory = "docs/ali"

for row in rows:
    docx_file = os.path.join(directory, row[0] + ".pptx")
    if os.path.exists(docx_file):
        process_pptx(docx_file)
        cur.execute(
            sql.SQL("UPDATE ali SET unstructure_time = %s WHERE id = %s"),
            [datetime.now(), row[0]],
        )
        conn.commit()
    else:
        continue

cur.close()
conn.close()

print("Data inserted successfully")
